from __future__ import annotations

import re
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent
SOURCE_MD = BASE_DIR / "presentation-offre-prospection-olithea-notebooklm.md"
OUTPUT_PPTX = BASE_DIR / "presentation-offre-prospection-olithea.pptx"
TEMP_ASSET_DIR = BASE_DIR / ".presentation-assets-tmp"

LOGO = PROJECT_DIR / "public/images/brand/olithea-logo-horizontal-accent-cropped.png"
MARK = PROJECT_DIR / "public/images/brand/olithea-mark-cropped.png"
SCREEN_CREATE = PROJECT_DIR / "public/images/features/parcours-offre-creation.webp"
SCREEN_MOBILE = PROJECT_DIR / "public/images/features/parcours-offre-capture-mobile.webp"
SCREEN_RESULTS = PROJECT_DIR / "public/images/features/parcours-offre-resultats.webp"

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "olive": "647A0B",
    "olive_dark": "526509",
    "dark": "2F3B24",
    "brown": "854F38",
    "brown_dark": "6B3E2D",
    "cream": "F7F9EC",
    "cream_2": "EEF3DD",
    "warm": "F4ECE7",
    "white": "FFFFFF",
    "text": "27302A",
    "muted": "667064",
    "line": "D8DECC",
    "gray": "E9ECE6",
    "black": "111511",
}

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = COLORS["text"],
    bold: bool = False,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
    line_spacing: float | None = None,
    uppercase: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_rich_text(slide, parts, x, y, w, h, *, size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for part in parts:
        run = p.add_run()
        run.text = part["text"]
        run.font.name = part.get("font", FONT_BODY)
        run.font.size = Pt(part.get("size", size))
        run.font.bold = part.get("bold", False)
        run.font.color.rgb = rgb(part.get("color", COLORS["text"]))
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = COLORS["white"],
    line: str | None = None,
    radius: bool = True,
    line_width: float = 1.0,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_shadow_card(slide, x, y, w, h, *, fill=COLORS["white"], line=None, offset=0.06):
    add_rect(slide, x + offset, y + offset, w, h, fill=COLORS["line"], line=None, radius=True)
    return add_rect(slide, x, y, w, h, fill=fill, line=line, radius=True)


def add_circle(slide, x, y, d, *, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=COLORS["line"], width=2.0):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    return connector


def add_pill(slide, text, x, y, w, h, *, fill, color=COLORS["white"], line=None, size=11, bold=True):
    add_rect(slide, x, y, w, h, fill=fill, line=line, radius=True)
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        uppercase=True,
    )


def add_kicker(slide, text, *, dark=False, x=0.72, y=0.46, w=6.0):
    color = COLORS["cream_2"] if dark else COLORS["olive"]
    return add_text(slide, text, x, y, w, 0.3, size=10.5, color=color, bold=True, uppercase=True)


def add_title(slide, text, *, dark=False, x=0.72, y=0.78, w=11.8, h=0.8, size=30):
    color = COLORS["white"] if dark else COLORS["dark"]
    return add_text(slide, text, x, y, w, h, size=size, color=color, bold=True, font=FONT_HEAD)


def add_footer(slide, slide_number: int, *, dark=False):
    color = "C7D0BB" if dark else "85907F"
    add_text(slide, "OLITHEA  •  OFFRE & PROSPECTION", 0.72, 7.15, 4.7, 0.18, size=8.5, color=color, bold=True)
    add_text(slide, f"{slide_number:02d}", 12.1, 7.12, 0.5, 0.22, size=9, color=color, bold=True, align=PP_ALIGN.RIGHT)


def add_logo(slide, x, y, w, h, *, card=False):
    if card:
        add_rect(slide, x - 0.15, y - 0.08, w + 0.3, h + 0.16, fill=COLORS["cream"], radius=True)
    return add_picture_contain(slide, LOGO, x, y, w, h)


def add_picture_contain(slide, path: Path, x, y, w, h):
    path = ensure_powerpoint_image(path)
    with Image.open(path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio >= box_ratio:
        pw = w
        ph = w / image_ratio
        px = x
        py = y + (h - ph) / 2
    else:
        ph = h
        pw = h * image_ratio
        px = x + (w - pw) / 2
        py = y
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def ensure_powerpoint_image(path: Path) -> Path:
    if path.suffix.lower() != ".webp":
        return path
    TEMP_ASSET_DIR.mkdir(parents=True, exist_ok=True)
    converted = TEMP_ASSET_DIR / f"{path.stem}.png"
    if not converted.exists() or converted.stat().st_mtime < path.stat().st_mtime:
        with Image.open(path) as source:
            source.convert("RGB").save(converted, format="PNG", optimize=True)
    return converted


def add_browser_frame(slide, path: Path, x, y, w, h, *, label=None):
    add_shadow_card(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"])
    add_rect(slide, x, y, w, 0.34, fill=COLORS["gray"], line=None, radius=True)
    for index, color in enumerate(("C56E5E", "D8A948", "7C9C52")):
        add_circle(slide, x + 0.16 + index * 0.19, y + 0.105, 0.09, fill=color)
    add_picture_contain(slide, path, x + 0.08, y + 0.43, w - 0.16, h - 0.52)
    if label:
        add_pill(slide, label, x + 0.3, y - 0.18, 2.3, 0.36, fill=COLORS["olive"], size=8.5)


def add_phone_frame(slide, path: Path, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=COLORS["black"], line=COLORS["black"], radius=True)
    add_picture_contain(slide, path, x + 0.10, y + 0.14, w - 0.20, h - 0.28)
    add_rect(slide, x + w * 0.34, y + 0.06, w * 0.32, 0.08, fill=COLORS["black"], radius=True)


def add_chevron(slide, x, y, w=0.34, h=0.52, *, color=COLORS["line"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_notes(slide, notes: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes.strip()


def markdown_to_notes(text: str) -> str:
    text = re.sub(r"(?m)^#{3,4}\s+", "\n", text)
    text = re.sub(r"(?m)^>\s?", "", text)
    text = text.replace("**", "").replace("`", "")
    text = re.sub(r"(?m)^-\s+", "• ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_notes() -> dict[int, str]:
    source = SOURCE_MD.read_text(encoding="utf-8")
    notes: dict[int, str] = {}
    pattern = re.compile(
        r"(?ms)^##\s+\d+\.\s+Diapositive\s+(\d+)\s+.*?\n(.*?)(?=^##\s+\d+\.\s+Diapositive\s+\d+|^---$)"
    )
    for match in pattern.finditer(source):
        slide_number = int(match.group(1))
        body = match.group(2)
        script_match = re.search(r"(?ms)^###\s+Script oral.*?\n(.*)$", body)
        if script_match:
            notes[slide_number] = markdown_to_notes(script_match.group(1))
    return notes


def slide_1(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["dark"])
    add_logo(slide, 0.85, 0.50, 2.15, 0.66, card=True)
    add_pill(slide, "20 MIN • POUR LES PRATICIENS", 0.82, 1.55, 2.75, 0.42, fill=COLORS["brown"], size=9)
    add_text(slide, "Une offre.\nDeux parcours.", 0.78, 2.05, 6.2, 1.85, size=42, color=COLORS["white"], bold=True, font=FONT_HEAD)
    add_text(slide, "De la découverte au rendez-vous", 0.82, 4.15, 5.7, 0.55, size=21, color=COLORS["cream_2"], bold=False)
    add_text(slide, "Offre & prospection • Olithea", 0.82, 5.04, 4.0, 0.35, size=11, color="BEC9AF", bold=True, uppercase=True)

    # Diagramme de couverture : intérêt -> deux rythmes -> rendez-vous.
    add_circle(slide, 9.45, 0.90, 1.20, fill=COLORS["cream"])
    add_text(slide, "INTÉRÊT", 9.45, 0.90, 1.20, 1.20, size=9, color=COLORS["dark"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_line(slide, 10.05, 2.10, 8.65, 2.85, color="97A98A", width=2.2)
    add_line(slide, 10.05, 2.10, 11.40, 2.85, color="B78A74", width=2.2)
    add_rect(slide, 7.75, 2.78, 2.15, 1.06, fill=COLORS["olive"], radius=True)
    add_text(slide, "CLAIRE", 7.96, 2.97, 1.72, 0.26, size=10, color=COLORS["white"], bold=True, uppercase=True)
    add_text(slide, "prête à agir", 7.96, 3.28, 1.72, 0.25, size=12.5, color=COLORS["white"], bold=True)
    add_rect(slide, 10.32, 2.78, 2.15, 1.06, fill=COLORS["brown"], radius=True)
    add_text(slide, "SOPHIE", 10.53, 2.97, 1.72, 0.26, size=10, color=COLORS["white"], bold=True, uppercase=True)
    add_text(slide, "veut découvrir", 10.53, 3.28, 1.72, 0.25, size=12.5, color=COLORS["white"], bold=True)
    add_line(slide, 8.82, 3.84, 10.05, 4.74, color="97A98A", width=2.2)
    add_line(slide, 11.40, 3.84, 10.05, 4.74, color="B78A74", width=2.2)
    add_circle(slide, 9.25, 4.65, 1.60, fill=COLORS["cream_2"], line="B9C5A4", line_width=1.2)
    add_text(slide, "RDV", 9.25, 4.65, 1.60, 1.60, size=19, color=COLORS["dark"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "La bonne prochaine étape\nau bon moment.", 8.05, 6.10, 4.15, 0.60, size=13.5, color=COLORS["cream_2"], bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 1, dark=True)
    add_notes(slide, notes.get(1, ""))


def slide_2(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["cream"])
    add_kicker(slide, "LE VRAI PROBLÈME")
    add_title(slide, "Mêmes praticienne et offre. Deux attentes.", size=29)
    add_text(slide, "Tout le monde ne se trouve pas au même niveau de décision.", 0.74, 1.35, 8.8, 0.4, size=14, color=COLORS["muted"])

    # Carte Claire
    add_shadow_card(slide, 0.75, 1.95, 5.75, 4.35, fill=COLORS["white"], line="CDD8BB")
    add_circle(slide, 1.15, 2.35, 0.95, fill=COLORS["olive"])
    add_text(slide, "C", 1.15, 2.35, 0.95, 0.95, size=23, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "CLAIRE", 2.35, 2.30, 2.5, 0.30, size=11, color=COLORS["olive"], bold=True, uppercase=True)
    add_text(slide, "Elle veut agir maintenant", 2.35, 2.64, 3.3, 0.45, size=21, color=COLORS["dark"], bold=True, font=FONT_HEAD)
    add_text(slide, "Recommandée par une amie", 1.20, 3.52, 4.2, 0.28, size=13, color=COLORS["text"], bold=True)
    add_text(slide, "Elle cherche le tarif, les modalités\net un créneau disponible.", 1.20, 3.94, 4.7, 0.80, size=15, color=COLORS["muted"])
    add_pill(slide, "PRENDRE RENDEZ-VOUS", 1.18, 5.35, 2.55, 0.50, fill=COLORS["olive"], size=9.5)

    # Carte Sophie
    add_shadow_card(slide, 6.82, 1.95, 5.75, 4.35, fill=COLORS["white"], line="DFC9BF")
    add_circle(slide, 7.22, 2.35, 0.95, fill=COLORS["brown"])
    add_text(slide, "S", 7.22, 2.35, 0.95, 0.95, size=23, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "SOPHIE", 8.42, 2.30, 2.5, 0.30, size=11, color=COLORS["brown"], bold=True, uppercase=True)
    add_text(slide, "Elle veut d’abord comprendre", 8.42, 2.64, 3.55, 0.55, size=21, color=COLORS["dark"], bold=True, font=FONT_HEAD)
    add_text(slide, "Arrivée depuis Instagram", 7.27, 3.52, 4.2, 0.28, size=13, color=COLORS["text"], bold=True)
    add_text(slide, "Elle se reconnaît, mais doit encore\ndécouvrir l’approche.", 7.27, 3.94, 4.7, 0.80, size=15, color=COLORS["muted"])
    add_pill(slide, "DÉCOUVRIR UNE RESSOURCE", 7.25, 5.35, 2.80, 0.50, fill=COLORS["brown"], size=9.5)

    add_text(slide, "La prospection, c’est proposer la bonne prochaine étape.", 1.35, 6.56, 10.65, 0.40, size=17, color=COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)
    add_notes(slide, notes.get(2, ""))


def slide_3(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_kicker(slide, "AVANT LE PARCOURS")
    add_title(slide, "Une offre claire donne une direction.", size=30)

    add_shadow_card(slide, 0.75, 1.62, 5.15, 4.95, fill=COLORS["dark"], line=None)
    add_text(slide, "VOTRE PHRASE", 1.12, 1.98, 2.1, 0.25, size=10, color="B9C7A9", bold=True, uppercase=True)
    add_rich_text(
        slide,
        [
            {"text": "J’accompagne ", "size": 20, "color": COLORS["white"], "bold": True, "font": FONT_HEAD},
            {"text": "[QUI]", "size": 20, "color": "C9D99B", "bold": True, "font": FONT_HEAD},
            {"text": "\nqui rencontre ", "size": 20, "color": COLORS["white"], "bold": True, "font": FONT_HEAD},
            {"text": "[SITUATION]", "size": 20, "color": "E0B3A0", "bold": True, "font": FONT_HEAD},
            {"text": "\nà aller vers ", "size": 20, "color": COLORS["white"], "bold": True, "font": FONT_HEAD},
            {"text": "[OBJECTIF]", "size": 20, "color": "C9D99B", "bold": True, "font": FONT_HEAD},
            {"text": "\ngrâce à ", "size": 20, "color": COLORS["white"], "bold": True, "font": FONT_HEAD},
            {"text": "[FORMAT]", "size": 20, "color": "E0B3A0", "bold": True, "font": FONT_HEAD},
            {"text": "\nen commençant par ", "size": 20, "color": COLORS["white"], "bold": True, "font": FONT_HEAD},
            {"text": "[ÉTAPE]", "size": 20, "color": "C9D99B", "bold": True, "font": FONT_HEAD},
        ],
        1.12,
        2.48,
        4.25,
        2.95,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(slide, "Sans promesse irréaliste. Sans jargon métier.", 1.12, 5.72, 4.15, 0.42, size=12.5, color="CDD8C1", bold=True)

    questions = [
        ("01", "POUR QUI ?", "Une personne doit se reconnaître."),
        ("02", "QUELLE SITUATION ?", "Un problème ou un moment précis."),
        ("03", "QUEL OBJECTIF ?", "Un progrès compréhensible et réaliste."),
        ("04", "SOUS QUELLE FORME ?", "Séance, atelier, programme…"),
        ("05", "QUEL PREMIER PAS ?", "Une seule prochaine action."),
    ]
    for idx, (number, title, detail) in enumerate(questions):
        y = 1.62 + idx * 0.96
        add_rect(slide, 6.38, y, 5.95, 0.78, fill=COLORS["cream"], line=COLORS["line"], radius=True)
        add_circle(slide, 6.60, y + 0.13, 0.52, fill=COLORS["olive"] if idx % 2 == 0 else COLORS["brown"])
        add_text(slide, number, 6.60, y + 0.13, 0.52, 0.52, size=9, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, 7.34, y + 0.13, 2.2, 0.25, size=10.5, color=COLORS["dark"], bold=True, uppercase=True)
        add_text(slide, detail, 7.34, y + 0.39, 4.45, 0.23, size=11.5, color=COLORS["muted"])

    add_text(slide, "Une prestation décrit ce que vous faites. Une offre aide la personne à se situer.", 6.40, 6.52, 5.85, 0.42, size=13, color=COLORS["brown"], bold=True)
    add_footer(slide, 3)
    add_notes(slide, notes.get(3, ""))


def slide_4(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["cream"])
    add_kicker(slide, "LE CHEMIN DE CLAIRE")
    add_title(slide, "Elle est prête : enlevez les détours.", size=30)
    add_text(slide, "Recommandée, décidée, pressée d’agir.", 0.74, 1.33, 6.4, 0.35, size=14, color=COLORS["muted"])

    nodes = [
        ("01", "ARRIVÉE", "Recommandation\nGoogle • Instagram"),
        ("02", "PORTAIL PUBLIC", "Qui • quoi • tarif\nmodalités • confiance"),
        ("03", "PRESTATION", "Un service clair\net un créneau"),
        ("04", "RENDEZ-VOUS", "Confirmation\nimmédiate"),
    ]
    start_x = 0.74
    for idx, (number, title, detail) in enumerate(nodes):
        x = start_x + idx * 3.08
        fill = COLORS["olive"] if idx in (1, 3) else COLORS["white"]
        text_color = COLORS["white"] if idx in (1, 3) else COLORS["dark"]
        detail_color = "E7EED8" if idx in (1, 3) else COLORS["muted"]
        add_shadow_card(slide, x, 2.02, 2.42, 2.05, fill=fill, line="C9D4B8" if idx not in (1, 3) else None)
        add_text(slide, number, x + 0.22, 2.23, 0.5, 0.25, size=10, color=detail_color, bold=True)
        add_text(slide, title, x + 0.22, 2.62, 1.98, 0.34, size=12, color=text_color, bold=True, uppercase=True)
        add_text(slide, detail, x + 0.22, 3.10, 1.98, 0.62, size=12.2, color=detail_color, bold=False)
        if idx < 3:
            add_chevron(slide, x + 2.60, 2.78, color=COLORS["brown"] if idx == 0 else COLORS["line"])

    add_rect(slide, 0.74, 4.62, 11.88, 1.82, fill=COLORS["dark"], line=None, radius=True)
    add_text(slide, "APRÈS LA RÉSERVATION, TOUT CONTINUE", 1.08, 4.93, 4.3, 0.30, size=10.5, color="C9D5B7", bold=True, uppercase=True)
    followups = ["Confirmation", "Rappels", "Agenda", "Note de séance", "Facture"]
    for idx, item in enumerate(followups):
        x = 1.06 + idx * 2.22
        add_pill(slide, item, x, 5.46, 1.88, 0.50, fill=COLORS["cream"], color=COLORS["dark"], size=9.5)
    add_text(slide, "Quand la personne est prête, le parcours ne doit pas la ralentir.", 0.78, 6.64, 11.70, 0.35, size=14.5, color=COLORS["olive_dark"], bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)
    add_notes(slide, notes.get(4, ""))


def slide_5(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_kicker(slide, "LE CHEMIN DE SOPHIE")
    add_title(slide, "Elle découvre : proposez un premier petit pas.", size=29)
    add_text(slide, "Intéressée, mais pas encore prête à réserver.", 0.74, 1.34, 6.8, 0.35, size=14, color=COLORS["muted"])

    steps = [
        ("01", "CONTENU", "Elle se reconnaît."),
        ("02", "RESSOURCE", "Elle reçoit une première valeur."),
        ("03", "SUIVI", "Elle comprend l’approche."),
        ("04", "OFFRE", "Elle voit la prochaine étape."),
        ("05", "RENDEZ-VOUS", "Elle agit lorsqu’elle est prête."),
    ]
    for idx, (number, title, detail) in enumerate(steps):
        y = 1.93 + idx * 0.88
        color = COLORS["brown"] if idx < 4 else COLORS["olive"]
        add_circle(slide, 0.86, y + 0.02, 0.58, fill=color)
        add_text(slide, number, 0.86, y + 0.02, 0.58, 0.58, size=9, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if idx < 4:
            add_line(slide, 1.15, y + 0.60, 1.15, y + 0.87, color="D8C2B7", width=2)
        add_text(slide, title, 1.70, y - 0.01, 1.80, 0.25, size=10.5, color=color, bold=True, uppercase=True)
        add_text(slide, detail, 3.22, y - 0.02, 3.65, 0.48, size=13.2, color=COLORS["text"], bold=True)

    add_rect(slide, 0.77, 6.35, 6.38, 0.58, fill=COLORS["warm"], line="DEC7BD", radius=True)
    add_text(slide, "Une étape utile avant une étape engageante.", 1.02, 6.35, 5.88, 0.58, size=14, color=COLORS["brown_dark"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Capture réelle dans une maquette de téléphone.
    add_rect(slide, 8.02, 1.38, 4.30, 5.72, fill=COLORS["cream"], line=COLORS["line"], radius=True)
    add_text(slide, "CÔTÉ VISITEUR", 8.35, 1.72, 1.75, 0.28, size=9, color=COLORS["brown"], bold=True, uppercase=True)
    add_text(slide, "Une page, une promesse,\nune seule action.", 8.35, 2.05, 2.20, 0.78, size=18, color=COLORS["dark"], bold=True, font=FONT_HEAD)
    add_phone_frame(slide, SCREEN_MOBILE, 10.30, 1.64, 1.58, 4.95)
    add_pill(slide, "RESSOURCE GRATUITE", 8.35, 3.30, 1.72, 0.42, fill=COLORS["brown"], size=8.5)
    add_text(slide, "Formulaire court\nConsentement clair\nProchaine étape visible", 8.35, 4.00, 1.72, 1.35, size=12.2, color=COLORS["muted"], bold=True)
    add_footer(slide, 5)
    add_notes(slide, notes.get(5, ""))


def slide_6(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["dark"])
    add_kicker(slide, "LE TUNNEL LE PLUS SIMPLE", dark=True)
    add_title(slide, "Une ressource. Trois messages. Une invitation.", dark=True, size=29)
    add_text(slide, "INFORMER  →  RASSURER  →  INVITER", 0.74, 1.42, 6.2, 0.35, size=11, color="C7D5B7", bold=True, uppercase=True)

    emails = [
        ("J0", "REMETTRE", "La ressource\n+ un petit résultat", COLORS["olive"]),
        ("J+2", "EXPLIQUER", "Une difficulté\n+ une réponse utile", COLORS["brown"]),
        ("J+4", "PROPOSER", "L’offre\n+ la prochaine étape", COLORS["olive"]),
    ]
    for idx, (day, title, detail, accent) in enumerate(emails):
        x = 0.77 + idx * 4.18
        add_shadow_card(slide, x, 2.05, 3.48, 2.50, fill=COLORS["cream"], line=None)
        add_pill(slide, day, x + 0.28, 2.34, 0.70, 0.42, fill=accent, size=9)
        add_text(slide, title, x + 1.18, 2.39, 1.75, 0.25, size=10.5, color=accent, bold=True, uppercase=True)
        add_text(slide, detail, x + 0.31, 3.02, 2.82, 0.90, size=18, color=COLORS["dark"], bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        if idx < 2:
            add_chevron(slide, x + 3.70, 2.95, 0.30, 0.54, color="8FA17E")

    add_rect(slide, 0.78, 5.04, 11.78, 1.37, fill=COLORS["brown"], line=None, radius=True)
    add_text(slide, "NEWSLETTER", 1.10, 5.32, 1.45, 0.25, size=10, color="F1DCD2", bold=True, uppercase=True)
    add_text(slide, "Rester présent, sans repartir de zéro.", 2.54, 5.24, 4.95, 0.42, size=20, color=COLORS["white"], bold=True, font=FONT_HEAD)
    add_text(slide, "1 situation  •  1 conseil  •  1 prochaine étape", 7.72, 5.31, 4.18, 0.30, size=11.5, color="F3DED5", bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "La ressource crée le contact. Le suivi crée la relation.", 0.86, 6.67, 11.55, 0.33, size=14.2, color="D3DEC5", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6, dark=True)
    add_notes(slide, notes.get(6, ""))


def slide_7(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["cream"])
    add_kicker(slide, "LA DIFFÉRENCE OLITHEA")
    add_title(slide, "Deux portes d’entrée. Un seul environnement.", size=29)

    # Branches à gauche.
    add_shadow_card(slide, 0.76, 1.78, 3.35, 1.48, fill=COLORS["white"], line="C7D5B3")
    add_pill(slide, "JE SUIS PRÊT", 1.02, 2.06, 1.30, 0.38, fill=COLORS["olive"], size=8.5)
    add_text(slide, "Portail → rendez-vous", 1.02, 2.56, 2.55, 0.32, size=15, color=COLORS["dark"], bold=True)
    add_shadow_card(slide, 0.76, 4.25, 3.35, 1.48, fill=COLORS["white"], line="DFC8BD")
    add_pill(slide, "JE DÉCOUVRE", 1.02, 4.53, 1.42, 0.38, fill=COLORS["brown"], size=8.5)
    add_text(slide, "Ressource → suivi → rendez-vous", 1.02, 5.03, 2.70, 0.44, size=14, color=COLORS["dark"], bold=True)

    # Lignes de convergence.
    add_line(slide, 4.12, 2.52, 5.14, 3.38, color=COLORS["olive"], width=3)
    add_line(slide, 4.12, 4.98, 5.14, 4.12, color=COLORS["brown"], width=3)

    # Centre Olithea.
    add_circle(slide, 5.00, 2.54, 2.38, fill=COLORS["dark"], line="AAB99A", line_width=1.5)
    add_picture_contain(slide, MARK, 5.57, 2.91, 1.23, 1.23)
    add_text(slide, "OLITHEA", 5.20, 4.20, 1.98, 0.28, size=10.5, color=COLORS["cream_2"], bold=True, align=PP_ALIGN.CENTER, uppercase=True)

    add_line(slide, 7.40, 3.73, 8.16, 3.73, color="91A080", width=3)
    add_chevron(slide, 7.88, 3.48, 0.30, 0.50, color="91A080")

    # Sorties métier.
    outputs = [
        ("Agenda", 8.40, 1.90, COLORS["olive"]),
        ("Contacts", 10.33, 1.90, COLORS["brown"]),
        ("Dossiers", 8.40, 3.28, COLORS["dark"]),
        ("Paiements", 10.33, 3.28, COLORS["olive"]),
        ("Factures", 9.37, 4.66, COLORS["brown"]),
    ]
    for label, x, y, fill in outputs:
        add_rect(slide, x, y, 1.65, 0.98, fill=fill, line=None, radius=True)
        add_text(slide, label, x, y, 1.65, 0.98, size=12.5, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_text(slide, "Le marketing et le quotidien métier ne sont plus séparés.", 1.08, 6.46, 11.18, 0.38, size=16, color=COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)
    add_notes(slide, notes.get(7, ""))


def slide_8(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["dark"])
    add_kicker(slide, "DÉMONSTRATION", dark=True)
    add_title(slide, "Deux parcours, dans le même produit.", dark=True, size=29)
    add_pill(slide, "DÉMO LIVE", 11.18, 0.50, 1.20, 0.40, fill=COLORS["brown"], size=9)

    add_browser_frame(slide, SCREEN_CREATE, 0.73, 1.58, 7.56, 4.91, label="CÔTÉ PRATICIEN")
    add_text(slide, "Créer et relier", 1.05, 6.55, 2.3, 0.30, size=13.5, color=COLORS["cream_2"], bold=True)
    add_text(slide, "Offre • ressource • rendez-vous", 3.20, 6.56, 4.45, 0.28, size=11.5, color="AFBDA2")

    add_rect(slide, 8.68, 1.40, 3.93, 5.25, fill=COLORS["cream"], line="66735F", radius=True)
    add_text(slide, "CÔTÉ VISITEUR", 8.98, 1.72, 1.75, 0.28, size=9, color=COLORS["brown"], bold=True, uppercase=True)
    add_text(slide, "Comprendre\net agir", 8.98, 2.08, 1.55, 0.85, size=20, color=COLORS["dark"], bold=True, font=FONT_HEAD)
    add_phone_frame(slide, SCREEN_MOBILE, 10.63, 1.70, 1.55, 4.72)
    add_text(slide, "Une page claire\nUn consentement\nUne prochaine étape", 8.98, 3.45, 1.50, 1.45, size=11.8, color=COLORS["muted"], bold=True)
    add_pill(slide, "PUIS LE MÊME AGENDA", 8.98, 5.45, 1.55, 0.43, fill=COLORS["olive"], size=7.8)
    add_footer(slide, 8, dark=True)
    add_notes(slide, notes.get(8, ""))


def slide_9(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_kicker(slide, "À FAIRE CETTE SEMAINE")
    add_title(slide, "Votre plan dans les 48 heures.", size=30)

    plan = [
        ("01", "Choisissez une seule offre"),
        ("02", "Écrivez une phrase claire"),
        ("03", "Créez deux prochaines étapes"),
        ("04", "Vérifiez votre lien Instagram"),
    ]
    for idx, (number, label) in enumerate(plan):
        y = 1.72 + idx * 1.07
        add_circle(slide, 0.82, y, 0.70, fill=COLORS["olive"] if idx < 2 else COLORS["brown"])
        add_text(slide, number, 0.82, y, 0.70, 0.70, size=10, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, 1.78, y + 0.08, 4.35, 0.50, size=18, color=COLORS["dark"], bold=True, font=FONT_HEAD)

    add_shadow_card(slide, 6.55, 1.52, 5.80, 4.84, fill=COLORS["cream"], line=COLORS["line"])
    add_text(slide, "VOTRE LIEN INSTAGRAM", 6.95, 1.90, 2.80, 0.30, size=10, color=COLORS["olive"], bold=True, uppercase=True)
    add_text(slide, "Deux choix suffisent.", 6.95, 2.34, 4.2, 0.48, size=23, color=COLORS["dark"], bold=True, font=FONT_HEAD)
    add_rect(slide, 6.95, 3.12, 4.98, 1.06, fill=COLORS["olive"], line=None, radius=True)
    add_text(slide, "JE SUIS PRÊT", 7.25, 3.34, 1.42, 0.22, size=9, color="DCE7C7", bold=True, uppercase=True)
    add_text(slide, "Prendre rendez-vous", 8.65, 3.30, 2.80, 0.32, size=15, color=COLORS["white"], bold=True, align=PP_ALIGN.RIGHT)
    add_rect(slide, 6.95, 4.46, 4.98, 1.06, fill=COLORS["warm"], line="DABEB1", radius=True)
    add_text(slide, "JE DÉCOUVRE", 7.25, 4.68, 1.50, 0.22, size=9, color=COLORS["brown"], bold=True, uppercase=True)
    add_text(slide, "Recevoir une ressource", 8.63, 4.64, 2.82, 0.32, size=15, color=COLORS["brown_dark"], bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "Une offre. Deux portes. Pas dix liens.", 7.18, 5.83, 4.60, 0.34, size=13.5, color=COLORS["muted"], bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, "Commencez petit. Rendez le chemin évident.", 0.84, 6.50, 5.42, 0.38, size=14.5, color=COLORS["brown"], bold=True)
    add_footer(slide, 9)
    add_notes(slide, notes.get(9, ""))


def slide_10(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["dark"])
    add_logo(slide, 0.88, 0.53, 2.05, 0.62, card=True)
    add_pill(slide, "CANEVA GRATUIT", 0.83, 1.55, 1.45, 0.40, fill=COLORS["brown"], size=8.5)
    add_text(slide, "Construisez vos\ndeux parcours\nen 15 minutes.", 0.82, 2.05, 6.30, 2.20, size=36, color=COLORS["white"], bold=True, font=FONT_HEAD)
    add_text(slide, "Scannez le QR code ou ouvrez le lien dans le chat.", 0.85, 4.67, 5.70, 0.50, size=16, color=COLORS["cream_2"])

    add_rect(slide, 8.77, 1.16, 3.58, 4.20, fill=COLORS["cream"], line="C5D0B5", radius=True)
    add_rect(slide, 9.18, 1.55, 2.74, 2.74, fill=COLORS["white"], line=COLORS["olive"], radius=False, line_width=2.0)
    add_text(slide, "QR CODE\nÀ INSÉRER", 9.18, 1.55, 2.74, 2.74, size=17, color=COLORS["olive_dark"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "[URL COURTE DE LA RESSOURCE]", 9.03, 4.57, 3.07, 0.34, size=10, color=COLORS["muted"], bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.84, 5.66, 11.51, 0.94, fill=COLORS["cream_2"], line=None, radius=True)
    add_rich_text(
        slide,
        [
            {"text": "Votre communication crée l’attention. ", "size": 13.5, "color": COLORS["dark"], "bold": True},
            {"text": "Votre offre donne une direction. ", "size": 13.5, "color": COLORS["brown"], "bold": True},
            {"text": "Olithea relie le reste.", "size": 13.5, "color": COLORS["olive_dark"], "bold": True},
        ],
        1.20,
        5.94,
        10.80,
        0.36,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10, dark=True)
    add_notes(slide, notes.get(10, ""))


def build_presentation() -> Presentation:
    notes = extract_notes()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "Une offre, deux parcours : de la découverte au rendez-vous"
    prs.core_properties.subject = "Intervention Olithea — offre et prospection pour praticiens"
    prs.core_properties.author = "Olithea"
    prs.core_properties.keywords = "Olithea, offre, prospection, praticiens, tunnel de vente, newsletter"

    slide_1(prs, notes)
    slide_2(prs, notes)
    slide_3(prs, notes)
    slide_4(prs, notes)
    slide_5(prs, notes)
    slide_6(prs, notes)
    slide_7(prs, notes)
    slide_8(prs, notes)
    slide_9(prs, notes)
    slide_10(prs, notes)
    return prs


def main() -> None:
    for path in (SOURCE_MD, LOGO, MARK, SCREEN_CREATE, SCREEN_MOBILE, SCREEN_RESULTS):
        if not path.exists():
            raise FileNotFoundError(path)
    try:
        presentation = build_presentation()
        presentation.save(OUTPUT_PPTX)
        print(f"Created {OUTPUT_PPTX} ({OUTPUT_PPTX.stat().st_size} bytes, {len(presentation.slides)} slides)")
    finally:
        shutil.rmtree(TEMP_ASSET_DIR, ignore_errors=True)


if __name__ == "__main__":
    main()
